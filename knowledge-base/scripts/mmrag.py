#!/usr/bin/env python3
"""
mmrag - Multimodal RAG Knowledge Base CLI

Ingest videos, images, audio, documents into a local ChromaDB vector database
using Gemini Embedding 2 for multimodal embeddings and Gemini Flash for
generating text descriptions of non-text media.

Usage:
    mmrag.py ingest <path> [<path>...] [--collection NAME]
    mmrag.py query <question> [--top-k N] [--threshold F] [--max-tokens N] [--collection NAME] [--json] [--full]
    mmrag.py status [--collection NAME]
    mmrag.py list [--collection NAME]
    mmrag.py collections
    mmrag.py delete <path> [--collection NAME]
    mmrag.py reset --confirm
"""

import argparse
import difflib
import hashlib
import json
import mimetypes
import os
import shutil
import subprocess
import sys
import time
from functools import cmp_to_key
from pathlib import Path

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
# cortextOS env-var overrides (set by kb-*.sh scripts)
MMRAG_DIR = Path(os.environ.get("MMRAG_DIR", str(Path.home() / ".mmrag")))
CONFIG_FILE = Path(os.environ.get("MMRAG_CONFIG", str(MMRAG_DIR / "config.json")))
CHROMADB_DIR = Path(os.environ.get("MMRAG_CHROMADB_DIR", str(MMRAG_DIR / "chromadb")))
MEDIA_DIR = MMRAG_DIR / "media"
LOG_DIR = MMRAG_DIR / "logs"

VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
DOC_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".go",
             ".rs", ".java", ".cpp", ".c", ".sh", ".yaml", ".yml", ".toml",
             ".html", ".css", ".sql", ".rb", ".swift", ".kt", ".r", ".lua"}
SUPPORTED_EXTS = VIDEO_EXTS | AUDIO_EXTS | IMAGE_EXTS | DOC_EXTS | TEXT_EXTS
# Directory names whose entire subtree is excluded from ingest.
IGNORE_DIR_PARTS = {
    ".trash", ".obsidian", ".git", ".cache", ".venv", "__pycache__",
    "node_modules", ".next", ".turbo", "dist", "build", "archive",
    "deprecated", "worktrees", ".claude", ".venv-synth", ".claude-flow", "contacts-backup",
    "site-packages", "dist-info",
}
# File extensions never worth embedding (diagram XML, lockfiles, etc.).
IGNORE_FILE_EXTS = {".drawio", ".lock", ".log", ".tmp", ".vcf"}

# Defaults
DEFAULT_TEXT_CHUNK_SIZE = 1500
DEFAULT_TEXT_CHUNK_OVERLAP = 200
DEFAULT_VIDEO_CHUNK_SECONDS = 60
DEFAULT_VIDEO_OVERLAP_SECONDS = 15
DEFAULT_AUDIO_CHUNK_SECONDS = 60
DEFAULT_AUDIO_OVERLAP_SECONDS = 10
DEFAULT_EMBEDDING_DIMENSIONS = 768
DEFAULT_SIMILARITY_THRESHOLD = 0.0  # return everything by default, let caller filter
DEFAULT_MAX_TOKENS = 0  # 0 = unlimited
DEFAULT_PREVIEW_CHARS = 300
DEFAULT_RECENCY_WEIGHT = 0.3


def _env_float(name, default):
    raw = os.environ.get(name, "").strip()
    if not raw:
        return default
    try:
        return float(raw)
    except ValueError:
        return default


RECENCY_TIE_BAND = _env_float("MMRAG_RECENCY_TIE_BAND", 0.05)
RECENCY_MAX_LIFT = _env_float("MMRAG_RECENCY_MAX_LIFT", 0.05)
AUTHORITATIVE_SIM = _env_float("MMRAG_AUTHORITATIVE_SIM", 0.7)
RECENCY_HALF_LIFE_DAYS = {
    "decision": 30,
    "note": 30,
    "reference": 90,
    "media": 90,
    "other": 90,
    "policy": 365,
}
DEFAULT_HALF_LIFE_DAYS = 90
NEUTRAL_DECAY = 0.5
DELETE_BATCH_SIZE = 500
GET_BATCH_SIZE = 5000
DEFAULT_TOP_DOCS = 5
MAX_FULL_DOC_BYTES = 200 * 1024
AUTO_INDEX_BEGIN = "<!-- BEGIN AUTO-INDEX (managed by mmrag reindex — do not edit below) -->"
AUTO_INDEX_END = "<!-- END AUTO-INDEX -->"
DEFAULT_WIKI_ROOT = Path.home() / "code" / "knowledge-sync" / "wiki"
DELIVER_LOW_SCORE = 0.55
DELIVER_AMBIGUITY_GAP = 0.02
DEFAULT_RECONCILE_COLLECTION = "shared-clearworksai"
DEFAULT_RECONCILE_ROOTS = (
    Path.home() / "code" / "knowledge-sync" / "wiki",
    Path.home() / "code" / "knowledge-sync" / "raw",
)

# Pricing (per 1M tokens)
EMBEDDING_PRICE_PER_M = 0.20
FLASH_INPUT_PRICE_PER_M = 0.15
FLASH_OUTPUT_PRICE_PER_M = 0.60

# Retry classifier for the Gemini generate_content call inside ingest_pdf.
# Module-level so a fault-injection test client can reference the same set.
TRANSIENT_HTTP_CODES = {429, 500, 503}
TRANSIENT_STATUS_NAMES = {"UNAVAILABLE", "RESOURCE_EXHAUSTED"}

USAGE_FILE = MMRAG_DIR / "usage.json"

# ---------------------------------------------------------------------------
# Usage Tracker
# ---------------------------------------------------------------------------
_tracker = None  # module-level, set by cmd_ingest/cmd_query


class UsageTracker:
    def __init__(self, operation="unknown"):
        self.session = {
            "embedding_tokens": 0,
            "embedding_calls": 0,
            "generation_input_tokens": 0,
            "generation_output_tokens": 0,
            "generation_calls": 0,
            "started_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            "operation": operation,
        }

    def track_embedding(self, content):
        self.session["embedding_calls"] += 1
        if isinstance(content, str):
            self.session["embedding_tokens"] += int(len(content.split()) * 1.3)
        elif isinstance(content, list):
            for part in content:
                if isinstance(part, str):
                    self.session["embedding_tokens"] += int(len(part.split()) * 1.3)
                else:
                    try:
                        self.session["embedding_tokens"] += max(256, len(part.data) // 4)
                    except Exception:
                        self.session["embedding_tokens"] += 256

    def track_generation(self, response):
        self.session["generation_calls"] += 1
        um = getattr(response, "usage_metadata", None)
        if um:
            self.session["generation_input_tokens"] += getattr(um, "prompt_token_count", 0) or 0
            self.session["generation_output_tokens"] += getattr(um, "candidates_token_count", 0) or 0

    def cost(self):
        emb = (self.session["embedding_tokens"] / 1_000_000) * EMBEDDING_PRICE_PER_M
        gen_in = (self.session["generation_input_tokens"] / 1_000_000) * FLASH_INPUT_PRICE_PER_M
        gen_out = (self.session["generation_output_tokens"] / 1_000_000) * FLASH_OUTPUT_PRICE_PER_M
        return {
            "embedding": round(emb, 6),
            "generation_input": round(gen_in, 6),
            "generation_output": round(gen_out, 6),
            "total": round(emb + gen_in + gen_out, 6),
        }

    def persist(self):
        MMRAG_DIR.mkdir(parents=True, exist_ok=True)
        self.session["finished_at"] = time.strftime("%Y-%m-%dT%H:%M:%S")
        self.session["cost"] = self.cost()

        data = {"cumulative": {}, "sessions": []}
        if USAGE_FILE.exists():
            try:
                with open(USAGE_FILE) as f:
                    data = json.load(f)
            except (json.JSONDecodeError, KeyError):
                data = {"cumulative": {}, "sessions": []}

        data.setdefault("sessions", []).append(self.session)

        c = data.get("cumulative", {})
        for key in ["embedding_tokens", "embedding_calls",
                     "generation_input_tokens", "generation_output_tokens",
                     "generation_calls"]:
            c[key] = c.get(key, 0) + self.session[key]

        c["total_cost"] = round(sum(
            s.get("cost", {}).get("total", 0) for s in data["sessions"]
        ), 6)
        data["cumulative"] = c

        with open(USAGE_FILE, "w") as f:
            json.dump(data, f, indent=2)

    def summary_line(self):
        c = self.cost()
        return (f"  Tokens: {self.session['embedding_tokens']:,} embedding, "
                f"{self.session['generation_input_tokens']:,} gen-input, "
                f"{self.session['generation_output_tokens']:,} gen-output | "
                f"Cost: ${c['total']:.4f}")


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
def load_config():
    if not CONFIG_FILE.exists():
        print("ERROR: Config not found. Run setup first:")
        print(f"  bash {Path(__file__).parent / 'setup.sh'}")
        sys.exit(1)
    with open(CONFIG_FILE) as f:
        return json.load(f)


def get_api_key(config):
    key = os.environ.get("GEMINI_API_KEY") or config.get("gemini_api_key")
    if not key:
        print("ERROR: No Gemini API key. Set GEMINI_API_KEY or run setup.")
        sys.exit(1)
    return key

# ---------------------------------------------------------------------------
# Gemini clients
# ---------------------------------------------------------------------------
def _load_factory(dotted_path):
    """Resolve a dotted import path to a callable.

    Accepts 'pkg.mod.attr' or 'pkg.mod:attr'. The colon form is unambiguous when
    the attribute name collides with a submodule name, so it is preferred.
    """
    if ":" in dotted_path:
        module_path, _, attr_path = dotted_path.partition(":")
    else:
        module_path, _, attr_path = dotted_path.rpartition(".")
    if not module_path or not attr_path:
        raise ValueError(
            f"MMRAG_GEMINI_CLIENT_FACTORY {dotted_path!r} must be 'module.attr' or 'module:attr'"
        )
    import importlib
    obj = importlib.import_module(module_path)
    for part in attr_path.split("."):
        obj = getattr(obj, part)
    if not callable(obj):
        raise TypeError(
            f"MMRAG_GEMINI_CLIENT_FACTORY {dotted_path!r} resolved to non-callable {type(obj).__name__}"
        )
    return obj


def get_genai_client(api_key):
    """Construct a Gemini client.

    Default returns google.genai.Client(api_key=api_key) — byte-identical to
    the prior behavior. To inject a fake client (e.g. for testing the retry
    loop in ingest_pdf), set the env-var MMRAG_GEMINI_CLIENT_FACTORY to a
    dotted import path of a callable taking (api_key) and returning an object
    with .models.generate_content / .models.embed_content compatible shape.
    See knowledge-base/scripts/_test_clients/fault_injection.py for a reference.
    """
    factory_path = os.environ.get("MMRAG_GEMINI_CLIENT_FACTORY")
    if factory_path:
        return _load_factory(factory_path)(api_key)
    from google import genai
    return genai.Client(api_key=api_key)


def _retry_generate_content(client, *, model, contents, backoffs=(5, 15, 45)):
    """Call client.models.generate_content with bounded retries on transient APIErrors.

    Retries on HTTP code in TRANSIENT_HTTP_CODES or status name in
    TRANSIENT_STATUS_NAMES; re-raises immediately on any other APIError (auth,
    malformed request, etc.); re-raises last_err after all attempts exhausted.

    backoffs is a tuple of sleep seconds between attempts. len(backoffs) is the
    attempt count. Tests pass (0, 0, 0) to skip sleeps.
    """
    from google.genai import errors as _genai_errors
    last_err = None
    for attempt, backoff in enumerate(backoffs, start=1):
        try:
            return client.models.generate_content(model=model, contents=contents)
        except _genai_errors.APIError as e:
            last_err = e
            is_transient = (e.code in TRANSIENT_HTTP_CODES) or (e.status in TRANSIENT_STATUS_NAMES)
            if not is_transient:
                raise
            if attempt < len(backoffs):
                print(f"    Transient error (HTTP {e.code} {e.status or ''}); retrying in {backoff}s (attempt {attempt}/{len(backoffs)})")
                time.sleep(backoff)
            else:
                print(f"    Exhausted retries on transient error: HTTP {e.code} {e.status or ''}")
    raise last_err if last_err else RuntimeError("retry loop completed without response or error")


def embed_content(client, config, content, task_type="RETRIEVAL_DOCUMENT"):
    """Embed content using Gemini Embedding 2. Content can be text string or list of Parts."""
    from google.genai import types
    result = client.models.embed_content(
        model=config.get("embedding_model", "gemini-embedding-2-preview"),
        contents=content,
        config=types.EmbedContentConfig(
            output_dimensionality=config.get("embedding_dimensions", DEFAULT_EMBEDDING_DIMENSIONS),
            task_type=task_type,
        ),
    )
    if _tracker:
        _tracker.track_embedding(content)
    return result.embeddings[0].values


def embed_multimodal(client, config, description_text, media_bytes, mime_type):
    """
    Option B embedding: combine text description + raw media into one embedding.
    This captures both semantic text meaning AND visual/audio content.
    """
    from google.genai import types
    contents = [
        description_text,
        types.Part.from_bytes(data=media_bytes, mime_type=mime_type),
    ]
    return embed_content(client, config, contents)


def embed_query(client, config, query_text):
    """Embed a query string for retrieval."""
    return embed_content(client, config, query_text, task_type="RETRIEVAL_QUERY")


def describe_media(client, config, file_path, media_type="video"):
    """Use Gemini Flash to generate a text description of media."""
    from google.genai import types

    mime = mimetypes.guess_type(str(file_path))[0] or "application/octet-stream"
    with open(file_path, "rb") as f:
        data = f.read()

    prompts = {
        "video": (
            "Provide a detailed description of this video. Include:\n"
            "1. What is being shown/demonstrated\n"
            "2. Any text visible on screen\n"
            "3. Key concepts or topics discussed\n"
            "4. A transcript of any spoken words\n"
            "5. Step-by-step actions if it's a tutorial\n"
            "Be thorough - this description will be used for search and retrieval."
        ),
        "image": (
            "Describe this image in detail. Include:\n"
            "1. What is shown in the image\n"
            "2. Any text visible in the image\n"
            "3. Key concepts or topics depicted\n"
            "4. Colors, layout, and composition\n"
            "Be thorough - this description will be used for search and retrieval."
        ),
        "audio": (
            "Transcribe and describe this audio. Include:\n"
            "1. A full transcript of spoken words\n"
            "2. Description of any sounds or music\n"
            "3. Key topics discussed\n"
            "4. Speaker identification if possible\n"
            "Be thorough - this description will be used for search and retrieval."
        ),
    }

    response = client.models.generate_content(
        model=config.get("gemini_model", "gemini-2.5-flash"),
        contents=[
            types.Part.from_bytes(data=data, mime_type=mime),
            prompts.get(media_type, prompts["video"]),
        ],
    )
    if _tracker:
        _tracker.track_generation(response)
    return response.text, data, mime

# ---------------------------------------------------------------------------
# ChromaDB
# ---------------------------------------------------------------------------
def get_chroma_collection(collection_name="default"):
    import chromadb
    client = chromadb.PersistentClient(path=str(CHROMADB_DIR))
    return client.get_or_create_collection(
        name=collection_name,
        metadata={"hnsw:space": "cosine"},
    )


def get_chroma_client():
    import chromadb
    return chromadb.PersistentClient(path=str(CHROMADB_DIR))

# ---------------------------------------------------------------------------
# Text chunking
# ---------------------------------------------------------------------------
def chunk_text(text, chunk_size=DEFAULT_TEXT_CHUNK_SIZE, overlap=DEFAULT_TEXT_CHUNK_OVERLAP):
    """Split text into overlapping chunks, preferring paragraph/section boundaries."""
    if len(text) <= chunk_size:
        return [text] if text.strip() else []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size

        # Try to break at a paragraph boundary (double newline)
        if end < len(text):
            # Look backwards from end for a good break point
            search_zone = text[max(start + chunk_size // 2, start):end]
            # Prefer double newline (paragraph break)
            para_break = search_zone.rfind("\n\n")
            if para_break != -1:
                end = max(start + chunk_size // 2, start) + para_break + 2
            else:
                # Fall back to single newline
                line_break = search_zone.rfind("\n")
                if line_break != -1:
                    end = max(start + chunk_size // 2, start) + line_break + 1

        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
        if start >= len(text):
            break

    return chunks

# ---------------------------------------------------------------------------
# Video chunking with FFmpeg
# ---------------------------------------------------------------------------
def get_media_duration(file_path):
    """Get duration of a media file in seconds. Returns 0 if unreadable."""
    try:
        result = subprocess.run(
            ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
             "-of", "csv=p=0", str(file_path)],
            capture_output=True, text=True,
        )
        val = result.stdout.strip()
        return float(val) if val else 0.0
    except (ValueError, subprocess.SubprocessError):
        return 0.0


def chunk_video(video_path, chunk_seconds=DEFAULT_VIDEO_CHUNK_SECONDS,
                overlap_seconds=DEFAULT_VIDEO_OVERLAP_SECONDS):
    """Split video into overlapping chunks using FFmpeg."""
    video_path = Path(video_path)
    output_dir = MEDIA_DIR / video_path.stem
    output_dir.mkdir(parents=True, exist_ok=True)

    duration = get_media_duration(video_path)

    chunks = []
    start = 0
    idx = 0
    step = chunk_seconds - overlap_seconds

    while start < duration:
        end = min(start + chunk_seconds, duration)
        # Skip tiny trailing chunks (< 5 seconds)
        if end - start < 5 and idx > 0:
            break

        output_file = output_dir / f"chunk_{idx:04d}.mp4"

        if not output_file.exists():
            subprocess.run(
                ["ffmpeg", "-y", "-i", str(video_path),
                 "-ss", str(start), "-t", str(end - start),
                 "-c", "copy", "-avoid_negative_ts", "1",
                 str(output_file)],
                capture_output=True,
            )

        chunks.append({
            "path": str(output_file),
            "start": start,
            "end": end,
            "index": idx,
        })

        start += step
        idx += 1

    return chunks


def chunk_audio(audio_path, chunk_seconds=DEFAULT_AUDIO_CHUNK_SECONDS,
                overlap_seconds=DEFAULT_AUDIO_OVERLAP_SECONDS):
    """Split audio into overlapping chunks using FFmpeg."""
    audio_path = Path(audio_path)
    output_dir = MEDIA_DIR / audio_path.stem
    output_dir.mkdir(parents=True, exist_ok=True)

    duration = get_media_duration(audio_path)

    ext = audio_path.suffix
    chunks = []
    start = 0
    idx = 0
    step = chunk_seconds - overlap_seconds

    while start < duration:
        end = min(start + chunk_seconds, duration)
        if end - start < 3 and idx > 0:
            break

        output_file = output_dir / f"chunk_{idx:04d}{ext}"

        if not output_file.exists():
            subprocess.run(
                ["ffmpeg", "-y", "-i", str(audio_path),
                 "-ss", str(start), "-t", str(end - start),
                 "-c", "copy", str(output_file)],
                capture_output=True,
            )

        chunks.append({
            "path": str(output_file),
            "start": start,
            "end": end,
            "index": idx,
        })

        start += step
        idx += 1

    return chunks

# ---------------------------------------------------------------------------
# File ID helper
# ---------------------------------------------------------------------------
def file_id(path, chunk_idx=None):
    """Generate a stable ID for a file or chunk."""
    h = hashlib.md5(str(path).encode()).hexdigest()[:12]
    if chunk_idx is not None:
        return f"{h}_chunk{chunk_idx}"
    return h


def _is_ignored(path: Path) -> bool:
    """True if any path component is an ignored dir, or the file ext is ignored."""
    for part in path.parts:
        part_lower = part.lower()
        if part_lower in IGNORE_DIR_PARTS:
            return True
        if part_lower.startswith(".venv"):
            return True
        if part_lower.endswith(".dist-info"):
            return True
    if path.suffix.lower() in IGNORE_FILE_EXTS:
        return True
    return False


def _normalize_source_path(file_path: Path) -> str:
    """Canonical absolute source path for stored metadata."""
    return str(Path(file_path).resolve())


def _source_mtime(file_path: Path) -> float:
    """mtime of the source file, or 0.0 if unavailable."""
    try:
        return Path(file_path).stat().st_mtime
    except OSError:
        return 0.0


def _source_created_at(file_path: Path) -> str:
    """ISO timestamp of the source file mtime, or '' if unavailable."""
    try:
        return time.strftime("%Y-%m-%dT%H:%M:%S", time.localtime(file_path.stat().st_mtime))
    except OSError:
        return ""


def _file_content_hash(file_path: Path) -> str:
    """sha256 of a source file's bytes."""
    digest = hashlib.sha256()
    with open(file_path, "rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _classify_doc_type(file_path: Path) -> str:
    """Coarse document class used to pick a recency half-life."""
    path_text = str(file_path).lower()
    if "/daily/" in path_text or "/sessions/" in path_text:
        return "note"
    if "/decisions/" in path_text or "decision" in path_text:
        return "decision"
    if "/wiki/" in path_text:
        return "reference"
    if any(key in path_text for key in ("policy", "/sop", "reference/", "definition", "glossary")):
        return "policy"
    if file_path.suffix.lower() in (IMAGE_EXTS | VIDEO_EXTS | AUDIO_EXTS):
        return "media"
    return "other"


def _common_source_metadata(file_path: Path, *, content_hash=None, source_mtime=None):
    """Additive metadata shared by every chunk of a given source file."""
    source_path = _normalize_source_path(file_path)
    return {
        "source": source_path,
        "source_file": source_path,
        "filename": file_path.name,
        "file_ext": file_path.suffix.lower(),
        "created_at": _source_created_at(file_path),
        "doc_type": _classify_doc_type(file_path),
        "content_hash": _file_content_hash(file_path) if content_hash is None else content_hash,
        "source_mtime": _source_mtime(file_path) if source_mtime is None else source_mtime,
    }


def _metadata_source(metadata):
    """Canonical source path from stored chunk metadata."""
    metadata = metadata or {}
    return metadata.get("source_file") or metadata.get("source") or ""


def _is_supported_file(file_path: Path) -> bool:
    """True when the file extension is one the ingest pipeline handles directly."""
    return file_path.suffix.lower() in SUPPORTED_EXTS


def _parse_reconcile_roots(roots_arg):
    """Parse comma-delimited roots, or fall back to the canonical knowledge-sync roots."""
    if not roots_arg:
        return [root.resolve() for root in DEFAULT_RECONCILE_ROOTS]
    roots = []
    for root in roots_arg.split(","):
        root = root.strip()
        if root:
            roots.append(Path(root).expanduser().resolve())
    return roots


def _read_text_file(file_path: Path, *, max_bytes=None):
    """Read a file as UTF-8 text with replacement, optionally truncating bytes."""
    data = Path(file_path).read_bytes()
    truncated = False
    if max_bytes is not None and len(data) > max_bytes:
        data = data[:max_bytes]
        truncated = True
    return data.decode("utf-8", errors="replace"), truncated


def _humanize_filename(stem: str) -> str:
    parts = stem.replace("-", " ").replace("_", " ").split()
    if not parts:
        return stem
    return " ".join(part.capitalize() for part in parts)


def _load_markdown_frontmatter(file_path: Path):
    """Parse simple YAML frontmatter keys from a markdown file."""
    if file_path.suffix.lower() != ".md":
        return {}
    try:
        text, _ = _read_text_file(file_path)
    except OSError:
        return {}
    lines = text.splitlines()
    if len(lines) < 3 or lines[0].strip() != "---":
        return {}

    metadata = {}
    for line in lines[1:]:
        stripped = line.strip()
        if stripped == "---":
            break
        if ":" not in line:
            continue
        key, value = line.split(":", 1)
        key = key.strip().lower()
        value = value.strip().strip('"').strip("'")
        if key:
            metadata[key] = value
    else:
        return {}
    return metadata


def _extract_markdown_h1(file_path: Path) -> str:
    """First markdown H1 after any frontmatter block."""
    if file_path.suffix.lower() != ".md":
        return ""
    try:
        text, _ = _read_text_file(file_path)
    except OSError:
        return ""

    lines = text.splitlines()
    start_idx = 0
    if lines and lines[0].strip() == "---":
        for idx, line in enumerate(lines[1:], start=1):
            if line.strip() == "---":
                start_idx = idx + 1
                break
    for line in lines[start_idx:]:
        stripped = line.strip()
        if stripped.startswith("# "):
            return stripped[2:].strip()
    return ""


def _document_title_and_summary(file_path: Path):
    """Best-effort title + one-line summary for doc-level output and indexes."""
    metadata = _load_markdown_frontmatter(file_path)
    title = metadata.get("title") or _extract_markdown_h1(file_path) or _humanize_filename(file_path.stem)
    summary = metadata.get("description") or metadata.get("summary") or ""
    summary = " ".join(summary.split())
    return title, summary


def _wiki_link_target(article_path: Path, index_dir: Path) -> str:
    """Shortest wiki-link target from an index file to a sibling article."""
    relative = article_path.relative_to(index_dir)
    return relative.with_suffix("").as_posix()


def _render_auto_index_block(index_path: Path, article_paths):
    """Render the managed auto-index section for one wiki directory."""
    lines = [AUTO_INDEX_BEGIN]
    for article_path in sorted(article_paths):
        title, summary = _document_title_and_summary(article_path)
        line = f"- [[{_wiki_link_target(article_path, index_path.parent)}|{title}]]"
        if summary:
            line += f": {summary}"
        lines.append(line)
    lines.append(AUTO_INDEX_END)
    return "\n".join(lines) + "\n"


def _apply_managed_index_block(existing_text: str, managed_block: str) -> str:
    """Replace or append the managed auto-index block while preserving manual prose."""
    begin_idx = existing_text.find(AUTO_INDEX_BEGIN)
    end_idx = existing_text.find(AUTO_INDEX_END)
    if begin_idx != -1 and end_idx != -1 and end_idx > begin_idx:
        end_idx += len(AUTO_INDEX_END)
        prefix = existing_text[:begin_idx].rstrip("\n")
        suffix = existing_text[end_idx:].lstrip("\n")
        parts = []
        if prefix:
            parts.append(prefix)
        parts.append(managed_block.rstrip("\n"))
        if suffix:
            parts.append(suffix.rstrip("\n"))
        return "\n\n".join(parts) + "\n"

    stripped = existing_text.rstrip("\n")
    if not stripped:
        return managed_block
    return stripped + "\n\n" + managed_block


def _iter_wiki_directories(root: Path):
    """Walk a wiki tree while pruning ignored directories."""
    root = Path(root)
    if not root.exists():
        return
    for current_root, dirnames, _ in os.walk(root):
        current = Path(current_root)
        dirnames[:] = [
            dirname for dirname in dirnames
            if not _is_ignored(current / dirname)
        ]
        yield current


def _reindex_indexes(root: Path, *, dry_run=False):
    """Regenerate managed sections for wiki _index.md files."""
    root = Path(root).expanduser().resolve()
    report = {
        "root": str(root),
        "dry_run": dry_run,
        "indexes_updated": 0,
        "indexes_created": 0,
        "files_linked": 0,
        "updated_paths": [],
        "diffs": [],
    }

    for directory in _iter_wiki_directories(root):
        article_paths = sorted(
            path for path in directory.iterdir()
            if path.is_file()
            and path.suffix.lower() == ".md"
            and path.name != "_index.md"
            and not _is_ignored(path)
        )
        index_path = directory / "_index.md"
        if not index_path.exists() and len(article_paths) < 2:
            continue

        existing_text = ""
        if index_path.exists():
            existing_text, _ = _read_text_file(index_path)
        managed_block = _render_auto_index_block(index_path, article_paths)
        updated_text = _apply_managed_index_block(existing_text, managed_block)
        report["files_linked"] += len(article_paths)

        if existing_text == updated_text and index_path.exists():
            continue

        if not index_path.exists():
            report["indexes_created"] += 1
        report["indexes_updated"] += 1
        report["updated_paths"].append(str(index_path))

        diff_text = "".join(
            difflib.unified_diff(
                existing_text.splitlines(True),
                updated_text.splitlines(True),
                fromfile=str(index_path),
                tofile=str(index_path),
            )
        )
        if diff_text:
            report["diffs"].append(diff_text)

        if not dry_run:
            index_path.write_text(updated_text, encoding="utf-8")

    return report


def _iter_reconcile_files(roots):
    """Yield canonical disk files that the reconciler should mirror into the index."""
    files = {}
    for root in roots:
        root_path = Path(root).expanduser()
        if not root_path.exists():
            continue
        if root_path.is_file():
            candidates = [root_path]
        else:
            candidates = [path for path in root_path.rglob("*") if path.is_file() and not path.name.startswith(".")]
        for candidate in candidates:
            if _is_ignored(candidate) or not _is_supported_file(candidate):
                continue
            resolved = candidate.resolve()
            files[str(resolved)] = resolved
    return [files[key] for key in sorted(files)]


def _collect_index_state(collection):
    """Group indexed chunks by source path, with their ids and recorded content hashes."""
    all_data = _get_all_metadatas(collection)
    ids = all_data.get("ids") or []
    metadatas = all_data.get("metadatas") or []
    state = {}
    for idx, doc_id in enumerate(ids):
        metadata = metadatas[idx] if idx < len(metadatas) else {}
        source = _metadata_source(metadata)
        if not source:
            continue
        entry = state.setdefault(source, {"ids": [], "hashes": set(), "chunk_count": 0})
        entry["ids"].append(doc_id)
        entry["hashes"].add((metadata or {}).get("content_hash", ""))
        entry["chunk_count"] += 1
    return state


def _get_all_metadatas(collection):
    """Fetch ids and metadatas in bounded pages to avoid large Chroma SQL variable limits."""
    ids = []
    metadatas = []
    offset = 0

    while True:
        page = collection.get(include=["metadatas"], limit=GET_BATCH_SIZE, offset=offset)
        page_ids = page.get("ids") or []
        if not page_ids:
            break
        ids.extend(page_ids)
        metadatas.extend(page.get("metadatas") or [])
        if len(page_ids) < GET_BATCH_SIZE:
            break
        offset += GET_BATCH_SIZE

    return {"ids": ids, "metadatas": metadatas}


def _delete_ids_in_batches(collection, ids, *, label="chunk batch"):
    """Delete ids in bounded batches so huge Chroma deletes cannot crash compaction."""
    ids = list(ids)
    if not ids:
        return {
            "batch_count": 0,
            "failed_batches": 0,
            "deleted_count": 0,
            "succeeded": True,
        }

    batch_count = 0
    failed_batches = 0
    deleted_count = 0
    for start in range(0, len(ids), DELETE_BATCH_SIZE):
        batch = ids[start:start + DELETE_BATCH_SIZE]
        batch_count += 1
        try:
            collection.delete(ids=batch)
            deleted_count += len(batch)
        except Exception as exc:
            failed_batches += 1
            print(f"  ERROR deleting {label} batch {batch_count}: {exc}")

    return {
        "batch_count": batch_count,
        "failed_batches": failed_batches,
        "deleted_count": deleted_count,
        "succeeded": failed_batches == 0,
    }


def _purge_collection(collection, *, dry_run=False):
    """Delete indexed chunks whose source is ignored or missing from disk."""
    index_state = _collect_index_state(collection)
    report = {
        "collection": getattr(collection, "name", "unknown"),
        "purged_files": 0,
        "purged_chunks": 0,
        "kept_files": 0,
        "kept_chunks": 0,
        "reasons": {
            "ignored": {"files": 0, "chunks": 0},
            "missing_from_disk": {"files": 0, "chunks": 0},
        },
        "delete_failures": {"files": 0, "chunks": 0, "batches": 0},
    }

    for source, entry in index_state.items():
        source_path = Path(source)
        reason = None
        if _is_ignored(source_path):
            reason = "ignored"
        elif not source_path.exists():
            reason = "missing_from_disk"

        if reason is None:
            report["kept_files"] += 1
            report["kept_chunks"] += entry["chunk_count"]
            continue

        if dry_run:
            report["purged_files"] += 1
            report["purged_chunks"] += entry["chunk_count"]
            report["reasons"][reason]["files"] += 1
            report["reasons"][reason]["chunks"] += entry["chunk_count"]
            continue

        delete_result = _delete_ids_in_batches(collection, entry["ids"], label=source)
        if delete_result["succeeded"]:
            report["purged_files"] += 1
            report["purged_chunks"] += delete_result["deleted_count"]
            report["reasons"][reason]["files"] += 1
            report["reasons"][reason]["chunks"] += delete_result["deleted_count"]
            continue

        report["kept_files"] += 1
        report["kept_chunks"] += entry["chunk_count"]
        report["delete_failures"]["files"] += 1
        report["delete_failures"]["chunks"] += entry["chunk_count"]
        report["delete_failures"]["batches"] += delete_result["failed_batches"]

    return report


def _scan_disk_state(roots):
    """Map canonical source paths on disk to their current content hashes."""
    state = {}
    for file_path in _iter_reconcile_files(roots):
        state[_normalize_source_path(file_path)] = {
            "path": file_path,
            "content_hash": _file_content_hash(file_path),
        }
    return state


def _reconcile_collection(client, config, collection, roots, *, dry_run=False):
    """Make a collection mirror the current disk state for the supplied roots."""
    disk_state = _scan_disk_state(roots)
    index_state = _collect_index_state(collection)

    new_sources = []
    changed_sources = []
    removed_sources = []
    ignored_sources = []
    unchanged_sources = []

    for source, entry in index_state.items():
        source_path = Path(source)
        if _is_ignored(source_path):
            ignored_sources.append(source)
            continue

        disk_entry = disk_state.get(source)
        if disk_entry is None:
            removed_sources.append(source)
            continue

        hashes = entry["hashes"]
        if len(hashes) != 1 or "" in hashes or next(iter(hashes)) != disk_entry["content_hash"]:
            changed_sources.append(source)
        else:
            unchanged_sources.append(source)

    for source in disk_state:
        if source not in index_state:
            new_sources.append(source)

    successful_changed_sources = []
    successful_removed_sources = []
    successful_ignored_sources = []
    purged_chunks = 0
    delete_failures = {"files": 0, "chunks": 0, "batches": 0}
    failed_paths = []
    delete_outcomes = {}
    delete_sources = changed_sources + removed_sources + ignored_sources
    for source in delete_sources:
        entry = index_state[source]
        if dry_run:
            delete_outcomes[source] = True
            purged_chunks += entry["chunk_count"]
            continue

        delete_result = _delete_ids_in_batches(collection, entry["ids"], label=source)
        delete_outcomes[source] = delete_result["succeeded"]
        if delete_result["succeeded"]:
            purged_chunks += delete_result["deleted_count"]
            continue

        delete_failures["files"] += 1
        delete_failures["chunks"] += entry["chunk_count"]
        delete_failures["batches"] += delete_result["failed_batches"]

    for source in changed_sources:
        if delete_outcomes.get(source):
            successful_changed_sources.append(source)
    for source in removed_sources:
        if delete_outcomes.get(source):
            successful_removed_sources.append(source)
    for source in ignored_sources:
        if delete_outcomes.get(source):
            successful_ignored_sources.append(source)

    new_chunks = 0
    if not dry_run:
        for source in sorted(new_sources + successful_changed_sources):
            file_path = disk_state[source]["path"]
            try:
                new_chunks += ingest_file(client, config, collection, file_path)
            except Exception as exc:
                error_message = " ".join(str(exc).split()) or "<no message>"
                print(f"  SKIP (error): {file_path} — {type(exc).__name__}: {error_message}")
                failed_paths.append(str(file_path))
                continue

    total_files_indexed_after = len(index_state) - len(removed_sources) - len(ignored_sources) + len(new_sources)
    if not dry_run:
        total_files_indexed_after = len(_collect_index_state(collection))

    return {
        "collection": getattr(collection, "name", "unknown"),
        "dry_run": dry_run,
        "new_files": len(new_sources),
        "new_chunks": new_chunks,
        "changed_files": len(changed_sources) if dry_run else len(successful_changed_sources),
        "removed_files": len(removed_sources) if dry_run else len(successful_removed_sources),
        "ignored_files": len(ignored_sources) if dry_run else len(successful_ignored_sources),
        "purged_chunks": purged_chunks,
        "unchanged_files": len(unchanged_sources),
        "total_files_on_disk": len(disk_state),
        "total_files_indexed_after": total_files_indexed_after,
        "delete_failures": delete_failures,
        "failed_files": len(failed_paths),
        "failed_paths": failed_paths,
        "roots": [str(Path(root)) for root in roots],
    }


def _emit_report(report, *, json_out=False):
    """Print reconcile or purge results in JSON or readable text form."""
    if json_out:
        print(json.dumps(report, indent=2))
        return

    if "collections" in report:
        print("Purge ignored/missing sources")
        print("=" * 40)
        print(f"Purged files: {report['purged_files']}")
        print(f"Purged chunks: {report['purged_chunks']}")
        print(f"Kept files: {report['kept_files']}")
        print(f"Kept chunks: {report['kept_chunks']}")
        for reason, counts in report["reasons"].items():
            print(f"  {reason}: {counts['files']} files, {counts['chunks']} chunks")
        return

    print("Reconcile summary")
    print("=" * 40)
    print(f"Collection: {report['collection']}")
    print(f"Dry run: {'yes' if report['dry_run'] else 'no'}")
    print(f"New files: {report['new_files']}")
    print(f"Changed files: {report['changed_files']}")
    print(f"Removed files: {report['removed_files']}")
    print(f"Ignored files: {report['ignored_files']}")
    print(f"Unchanged files: {report['unchanged_files']}")
    print(f"Failed (skipped on error): {report.get('failed_files', 0)}")
    print(f"Purged chunks: {report['purged_chunks']}")
    print(f"New chunks: {report['new_chunks']}")
    print(f"Files on disk: {report['total_files_on_disk']}")
    print(f"Files indexed after: {report['total_files_indexed_after']}")
    if report.get("delete_failures", {}).get("files"):
        print("Delete failures:")
        print(f"  Files: {report['delete_failures']['files']}")
        print(f"  Chunks: {report['delete_failures']['chunks']}")
        print(f"  Batches: {report['delete_failures']['batches']}")
    if report.get("orphan_reap"):
        orphan_report = report["orphan_reap"]
        print()
        print("Orphan reaper")
        print("=" * 40)
        print(f"Collections scanned: {orphan_report['collections_scanned']}")
        print(f"Orphans found: {orphan_report['orphans_found']}")
        print(f"Orphans reaped: {orphan_report['orphans_reaped']}")
        print(f"Kept: {orphan_report['kept']}")
    if report.get("reindex"):
        reindex_report = report["reindex"]
        print()
        print("Auto-index refresh")
        print("=" * 40)
        print(f"Indexes updated: {reindex_report['indexes_updated']}")
        print(f"Indexes created: {reindex_report['indexes_created']}")
        print(f"Files linked: {reindex_report['files_linked']}")
        if reindex_report["dry_run"] and reindex_report["diffs"]:
            for diff_text in reindex_report["diffs"]:
                print(diff_text, end="" if diff_text.endswith("\n") else "\n")
    if report.get("failed_paths"):
        print()
        print("Failed paths")
        print("=" * 40)
        for failed_path in report["failed_paths"]:
            print(f"  {failed_path}")


def _recency_decay(created_at: str, doc_type: str, now: float = None) -> float:
    """0.5^(age_days / half_life), or NEUTRAL_DECAY when metadata is missing."""
    if not created_at:
        return NEUTRAL_DECAY
    try:
        created_ts = time.mktime(time.strptime(created_at, "%Y-%m-%dT%H:%M:%S"))
    except (ValueError, OverflowError):
        return NEUTRAL_DECAY
    current_time = time.time() if now is None else now
    age_days = max(0.0, (current_time - created_ts) / 86400.0)
    half_life = RECENCY_HALF_LIFE_DAYS.get(doc_type, DEFAULT_HALF_LIFE_DAYS)
    return 0.5 ** (age_days / half_life)


def _effective_recency_lift(recency: float, *, weight=DEFAULT_RECENCY_WEIGHT) -> float:
    """Bound recency to a tie-breaker sized score adjustment around neutral decay."""
    raw_lift = weight * (recency - NEUTRAL_DECAY)
    return max(-RECENCY_MAX_LIFT, min(RECENCY_MAX_LIFT, raw_lift))


def _compare_reranked_results(left, right):
    """Comparator that lets recency break ties without overriding strong relevance."""
    left_similarity = left["similarity"]
    right_similarity = right["similarity"]
    left_authoritative = left_similarity >= AUTHORITATIVE_SIM
    right_authoritative = right_similarity >= AUTHORITATIVE_SIM
    if left_authoritative != right_authoritative:
        return -1 if left_authoritative else 1

    similarity_gap = left_similarity - right_similarity
    if abs(similarity_gap) > RECENCY_TIE_BAND:
        return -1 if similarity_gap > 0 else 1

    score_gap = left["final_score"] - right["final_score"]
    if score_gap:
        return -1 if score_gap > 0 else 1
    if similarity_gap:
        return -1 if similarity_gap > 0 else 1
    return 0


def _apply_recency_rerank(results, *, enabled=True, weight=DEFAULT_RECENCY_WEIGHT, now=None):
    """Rescore post-threshold results with recency, or return them unchanged."""
    if not enabled:
        return results

    reranked = []
    for result in results:
        metadata = result.get("metadata") or {}
        recency = _recency_decay(
            metadata.get("created_at", ""),
            metadata.get("doc_type", "other"),
            now=now,
        )
        lift = _effective_recency_lift(recency, weight=weight)
        reranked.append({
            **result,
            "recency": recency,
            "final_score": result["similarity"] + lift,
        })
    reranked.sort(key=cmp_to_key(_compare_reranked_results))
    return reranked

# ---------------------------------------------------------------------------
# Ingest logic
# ---------------------------------------------------------------------------
def already_exists(collection, doc_id):
    """Check if a document ID already exists in the collection. Respects --force flag."""
    if args_force:
        return False
    existing = collection.get(ids=[doc_id])
    return bool(existing and existing["ids"])


def ingest_text_file(client, config, collection, file_path):
    """Ingest a text-based file."""
    file_path = Path(file_path)
    text = file_path.read_text(errors="replace")
    if not text.strip():
        print(f"  SKIP (empty): {file_path}")
        return 0
    common_metadata = _common_source_metadata(file_path)

    chunks = chunk_text(
        text,
        chunk_size=config.get("text_chunk_size", DEFAULT_TEXT_CHUNK_SIZE),
        overlap=config.get("text_chunk_overlap", DEFAULT_TEXT_CHUNK_OVERLAP),
    )

    count = 0
    for i, chunk in enumerate(chunks):
        doc_id = file_id(file_path, i)
        if already_exists(collection, doc_id):
            continue

        embedding = embed_content(client, config, chunk)
        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[chunk],
            metadatas=[{
                **common_metadata,
                "type": "text",
                "chunk_index": i,
                "total_chunks": len(chunks),
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        count += 1

    return count


def ingest_image(client, config, collection, file_path):
    """Ingest an image: Gemini Flash describes it, then embed description + raw image together."""
    file_path = Path(file_path)
    doc_id = file_id(file_path)
    common_metadata = _common_source_metadata(file_path)

    if already_exists(collection, doc_id):
        print(f"  SKIP (exists): {file_path}")
        return 0

    print(f"  Generating description for {file_path.name}...")
    description, media_bytes, mime = describe_media(client, config, file_path, "image")

    # Option B: embed text description + raw image together
    try:
        embedding = embed_multimodal(client, config, description, media_bytes, mime)
    except Exception:
        # Fallback to text-only embedding if multimodal fails (e.g., file too large)
        embedding = embed_content(client, config, description)

    collection.upsert(
        ids=[doc_id],
        embeddings=[embedding],
        documents=[description],
        metadatas=[{
            **common_metadata,
            "type": "image",
            "mime_type": mime,
            "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }],
    )
    return 1


def extract_audio_from_video(video_path, output_path):
    """Extract audio track from a video file as mp3."""
    subprocess.run(
        ["ffmpeg", "-y", "-i", str(video_path),
         "-vn", "-acodec", "libmp3lame", "-q:a", "4",
         str(output_path)],
        capture_output=True,
    )
    return Path(output_path).exists()


def ingest_video(client, config, collection, file_path):
    """Ingest a video: chunk it, describe each chunk, embed.

    For large videos (chunks > 20MB), falls back to audio-only extraction
    since the video bytes would be too large for the embedding API.
    For small chunks, uses full multimodal embedding (description + video).
    """
    file_path = Path(file_path)
    size_mb = file_path.stat().st_size / (1024 * 1024)
    duration = get_media_duration(file_path)

    if duration <= 0:
        print(f"  SKIP (unreadable/zero duration): {file_path}")
        return 0

    print(f"  Video: {file_path.name} ({size_mb:.0f}MB, {duration:.0f}s)")

    # Chunk the video
    chunk_secs = config.get("video_chunk_seconds", DEFAULT_VIDEO_CHUNK_SECONDS)
    overlap_secs = config.get("video_overlap_seconds", DEFAULT_VIDEO_OVERLAP_SECONDS)

    print(f"  Chunking into {chunk_secs}s segments with {overlap_secs}s overlap...")
    chunks = chunk_video(file_path, chunk_seconds=chunk_secs, overlap_seconds=overlap_secs)
    total_chunks = len(chunks)
    print(f"  Created {total_chunks} chunks")
    common_metadata = _common_source_metadata(file_path)

    count = 0
    for chunk in chunks:
        doc_id = file_id(file_path, chunk["index"])
        if already_exists(collection, doc_id):
            continue

        chunk_path = Path(chunk["path"])
        chunk_size_mb = chunk_path.stat().st_size / (1024 * 1024) if chunk_path.exists() else 0

        print(f"  Chunk {chunk['index'] + 1}/{total_chunks} "
              f"({chunk['start']:.0f}s-{chunk['end']:.0f}s, {chunk_size_mb:.1f}MB)")

        description = None
        media_bytes = None
        mime = None

        # Strategy: try video description first, fall back to audio-only for large chunks
        if chunk_size_mb <= 20:
            # Small enough for full video analysis
            try:
                description, media_bytes, mime = describe_media(client, config, chunk_path, "video")
                print(f"    Described via video")
            except Exception as e:
                print(f"    Video description failed ({e}), trying audio...")

        if description is None:
            # Large chunk or video failed: extract audio and describe that
            audio_path = chunk_path.with_suffix(".mp3")
            if not audio_path.exists():
                print(f"    Extracting audio track...")
                extract_audio_from_video(chunk_path, audio_path)

            if audio_path.exists() and audio_path.stat().st_size > 0:
                try:
                    description, media_bytes, mime = describe_media(client, config, audio_path, "audio")
                    # Prefix so the agent knows this came from a video's audio
                    description = (
                        f"[Audio extracted from video: {file_path.name}, "
                        f"{chunk['start']:.0f}s-{chunk['end']:.0f}s]\n\n{description}"
                    )
                    print(f"    Described via audio extraction")
                except Exception as e:
                    print(f"    Audio description also failed: {e}")

        if description is None:
            description = (
                f"Video chunk from {file_path.name}, "
                f"{chunk['start']:.0f}s to {chunk['end']:.0f}s. "
                f"(Description unavailable - file may be too large or corrupted)"
            )

        # Embed: try multimodal if we have small media bytes, else text-only
        if media_bytes and mime and len(media_bytes) < 20 * 1024 * 1024:
            try:
                embedding = embed_multimodal(client, config, description, media_bytes, mime)
            except Exception:
                embedding = embed_content(client, config, description)
        else:
            embedding = embed_content(client, config, description)

        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[description],
            metadatas=[{
                **common_metadata,
                "type": "video_chunk",
                "chunk_index": chunk["index"],
                "total_chunks": total_chunks,
                "chunk_start_seconds": chunk["start"],
                "chunk_end_seconds": chunk["end"],
                "chunk_path": str(chunk_path),
                "duration_seconds": duration,
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        count += 1

    return count


def ingest_audio(client, config, collection, file_path):
    """Ingest audio: chunk if needed, describe, embed description + audio together."""
    file_path = Path(file_path)
    duration = get_media_duration(file_path)
    common_metadata = _common_source_metadata(file_path)
    if duration <= 0:
        print(f"  SKIP (unreadable/zero duration): {file_path}")
        return 0
    max_chunk = config.get("audio_chunk_seconds", DEFAULT_AUDIO_CHUNK_SECONDS)

    if duration <= max_chunk:
        # Short enough to process as one piece
        doc_id = file_id(file_path)
        if already_exists(collection, doc_id):
            print(f"  SKIP (exists): {file_path}")
            return 0

        print(f"  Transcribing {file_path.name}...")
        description, media_bytes, mime = describe_media(client, config, file_path, "audio")

        try:
            embedding = embed_multimodal(client, config, description, media_bytes, mime)
        except Exception:
            embedding = embed_content(client, config, description)

        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[description],
            metadatas=[{
                **common_metadata,
                "type": "audio",
                "duration_seconds": duration,
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        return 1
    else:
        # Chunk the audio
        print(f"  Chunking audio: {file_path.name} ({duration:.0f}s)...")
        overlap = config.get("audio_overlap_seconds", DEFAULT_AUDIO_OVERLAP_SECONDS)
        chunks = chunk_audio(file_path, chunk_seconds=max_chunk, overlap_seconds=overlap)
        total_chunks = len(chunks)
        count = 0

        for chunk in chunks:
            doc_id = file_id(file_path, chunk["index"])
            if already_exists(collection, doc_id):
                continue

            print(f"  Transcribing chunk {chunk['index'] + 1}/{total_chunks}...")
            try:
                description, media_bytes, mime = describe_media(client, config, chunk["path"], "audio")
            except Exception as e:
                print(f"  WARNING: Failed to transcribe chunk {chunk['index']}: {e}")
                description = f"Audio chunk from {file_path.name}, {chunk['start']:.0f}s to {chunk['end']:.0f}s"
                media_bytes = None
                mime = None

            if media_bytes and mime:
                try:
                    embedding = embed_multimodal(client, config, description, media_bytes, mime)
                except Exception:
                    embedding = embed_content(client, config, description)
            else:
                embedding = embed_content(client, config, description)

            collection.upsert(
                ids=[doc_id],
                embeddings=[embedding],
                documents=[description],
                metadatas=[{
                    **common_metadata,
                    "type": "audio_chunk",
                    "chunk_index": chunk["index"],
                    "total_chunks": total_chunks,
                    "chunk_start_seconds": chunk["start"],
                    "chunk_end_seconds": chunk["end"],
                    "chunk_path": chunk["path"],
                    "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
                }],
            )
            count += 1
        return count


def ingest_pdf(client, config, collection, file_path):
    """Ingest a PDF page-by-page using Gemini to extract content including visual elements."""
    file_path = Path(file_path)
    if os.path.getsize(file_path) == 0:
        print(f"  SKIP (empty): {file_path}")
        return 0
    from google.genai import types
    common_metadata = _common_source_metadata(file_path)

    with open(file_path, "rb") as f:
        data = f.read()

    # Estimate page count (rough: ~3KB per page for typical PDFs, but varies wildly)
    # We'll ask Gemini to process the whole thing and get structured output
    # For PDFs > 6 pages, we chunk by asking for specific page ranges

    print(f"  Analyzing PDF: {file_path.name}...")

    # Gemini Flash returns 503 UNAVAILABLE during high-demand windows. Without
    # retries, a single 503 kills the ingest. _retry_generate_content wraps the
    # call with bounded retries on transient SDK conditions (HTTP 429/500/503,
    # status UNAVAILABLE/RESOURCE_EXHAUSTED) and fails fast on everything else.
    extraction_prompt = (
        "Extract ALL content from this PDF. For each page, include:\n"
        "1. Page number\n"
        "2. All text content (headings, body, lists, footnotes)\n"
        "3. Description of any images, charts, diagrams, or tables\n"
        "4. Key concepts and topics on that page\n"
        "Separate each page's content with '=== PAGE N ===' markers.\n"
        "Be thorough - this will be used for search and retrieval."
    )
    response = _retry_generate_content(
        client,
        model=config.get("gemini_model", "gemini-2.5-flash"),
        contents=[
            types.Part.from_bytes(data=data, mime_type="application/pdf"),
            extraction_prompt,
        ],
    )
    if _tracker:
        _tracker.track_generation(response)
    text = response.text

    # Split by page markers if present, otherwise chunk normally
    pages = []
    if "=== PAGE" in text:
        import re
        page_splits = re.split(r'===\s*PAGE\s*\d+\s*===', text)
        pages = [p.strip() for p in page_splits if p.strip()]
    else:
        # No page markers - chunk as text
        pages = chunk_text(
            text,
            chunk_size=config.get("text_chunk_size", DEFAULT_TEXT_CHUNK_SIZE),
            overlap=config.get("text_chunk_overlap", DEFAULT_TEXT_CHUNK_OVERLAP),
        )

    count = 0
    for i, page_content in enumerate(pages):
        if not page_content.strip():
            continue
        doc_id = file_id(file_path, i)
        if already_exists(collection, doc_id):
            continue

        embedding = embed_content(client, config, page_content)
        collection.upsert(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[page_content],
            metadatas=[{
                **common_metadata,
                "type": "pdf_page",
                "chunk_index": i,
                "total_chunks": len(pages),
                "page_number": i + 1,
                "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }],
        )
        count += 1
    return count


def extract_docx_text(file_path):
    """Extract text from .docx using python-docx."""
    from docx import Document
    doc = Document(str(file_path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading ", "").strip()
                prefix = "#" * (int(level) if level.isdigit() else 1)
                parts.append(f"{prefix} {para.text}")
            else:
                parts.append(para.text)
    # Also extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def extract_pptx_text(file_path):
    """Extract text from .pptx using python-pptx."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [f"=== SLIDE {i+1} ==="]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"Speaker Notes: {notes}")
        slides.append("\n".join(texts))
    return "\n\n".join(slides)


def extract_xlsx_text(file_path):
    """Extract text from .xlsx using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(str(file_path), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"=== SHEET: {sheet_name} ===\n" + "\n".join(rows[:200]))  # cap at 200 rows
    return "\n\n".join(parts)


def ingest_office_doc(client, config, collection, file_path):
    """Ingest Office documents (.docx, .pptx, .xlsx) by extracting text locally."""
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    common_metadata = _common_source_metadata(file_path)

    print(f"  Extracting content from {file_path.name}...")

    try:
        if ext in (".docx", ".doc"):
            text = extract_docx_text(file_path)
            type_name = "docx"
        elif ext in (".pptx", ".ppt"):
            text = extract_pptx_text(file_path)
            type_name = "slides"
        elif ext in (".xlsx", ".xls"):
            text = extract_xlsx_text(file_path)
            type_name = "spreadsheet"
        else:
            print(f"  SKIP (unsupported office format): {file_path}")
            return 0
    except Exception as e:
        print(f"  ERROR extracting {file_path.name}: {e}")
        return 0

    if not text.strip():
        print(f"  SKIP (empty document): {file_path}")
        return 0

    # Split presentations by slide markers, everything else by text chunks
    sections = []
    if type_name == "slides" and "=== SLIDE" in text:
        import re
        slide_splits = re.split(r'===\s*SLIDE\s*\d+\s*===', text)
        sections = [s.strip() for s in slide_splits if s.strip()]
    elif type_name == "spreadsheet" and "=== SHEET" in text:
        import re
        sheet_splits = re.split(r'===\s*SHEET:.*?===', text)
        sections = [s.strip() for s in sheet_splits if s.strip()]
    else:
        sections = chunk_text(
            text,
            chunk_size=config.get("text_chunk_size", DEFAULT_TEXT_CHUNK_SIZE),
            overlap=config.get("text_chunk_overlap", DEFAULT_TEXT_CHUNK_OVERLAP),
        )

    count = 0
    for i, section in enumerate(sections):
        if not section.strip():
            continue
        doc_id = file_id(file_path, i)
        if already_exists(collection, doc_id):
            continue

        embedding = embed_content(client, config, section)
        meta = {
            **common_metadata,
            "type": type_name,
            "chunk_index": i,
            "total_chunks": len(sections),
            "ingested_at": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        if type_name == "slides":
            meta["slide_number"] = i + 1

        collection.upsert(ids=[doc_id], embeddings=[embedding], documents=[section], metadatas=[meta])
        count += 1
    return count


# Global flag for --force re-ingestion
args_force = False


def ingest_file(client, config, collection, file_path):
    """Route a file to the appropriate ingest handler."""
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    # Skip common non-content files
    skip_names = {".ds_store", "thumbs.db", ".gitignore", ".gitkeep", "package-lock.json",
                  "yarn.lock", "pnpm-lock.yaml", ".eslintcache"}
    if file_path.name.lower() in skip_names:
        return 0

    # Skip junk directories
    skip_dirs = {".git", "node_modules", "__pycache__", ".venv", "venv", ".env",
                 ".next", ".nuxt", "dist", "build", ".cache", ".turbo",
                 "vendor", ".terraform", ".angular", ".svelte-kit", ".output",
                 "coverage", ".nyc_output", ".pytest_cache", ".mypy_cache"}
    parts = set(file_path.parts)
    if parts & skip_dirs:
        return 0

    size_bytes = os.path.getsize(file_path)
    if size_bytes == 0:
        print(f"  SKIP (empty): {file_path}")
        return 0

    # Skip text files > 10MB (likely generated/binary)
    size_mb = size_bytes / (1024 * 1024)
    if ext in TEXT_EXTS and size_mb > 10:
        print(f"  SKIP (too large: {size_mb:.0f}MB): {file_path}")
        return 0
    if ext in IMAGE_EXTS and size_mb > 50:
        print(f"  SKIP (too large: {size_mb:.0f}MB): {file_path}")
        return 0
    if ext in DOC_EXTS and size_mb > 100:
        print(f"  SKIP (too large: {size_mb:.0f}MB): {file_path}")
        return 0

    if ext in VIDEO_EXTS:
        return ingest_video(client, config, collection, file_path)
    elif ext in AUDIO_EXTS:
        return ingest_audio(client, config, collection, file_path)
    elif ext in IMAGE_EXTS:
        return ingest_image(client, config, collection, file_path)
    elif ext == ".pdf":
        return ingest_pdf(client, config, collection, file_path)
    elif ext in DOC_EXTS:
        return ingest_office_doc(client, config, collection, file_path)
    elif ext in TEXT_EXTS:
        return ingest_text_file(client, config, collection, file_path)
    else:
        # Try as text for unknown extensions
        try:
            file_path.read_text(errors="strict")[:100]
            return ingest_text_file(client, config, collection, file_path)
        except (UnicodeDecodeError, Exception):
            print(f"  SKIP (binary/unsupported): {file_path}")
            return 0

# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------
def cmd_ingest(args):
    global args_force, _tracker
    args_force = getattr(args, 'force', False)
    _tracker = UsageTracker("ingest")

    config = load_config()
    client = get_genai_client(get_api_key(config))
    collection_name = args.collection or config.get("default_collection", "default")
    collection = get_chroma_collection(collection_name)

    if args_force:
        print(f"Force mode: will re-ingest existing files")

    total = 0
    skipped = 0
    errors = 0

    try:
        for path_str in args.paths:
            p = Path(path_str).resolve()
            if p.is_dir():
                all_files = sorted(f for f in p.rglob("*") if f.is_file() and not f.name.startswith("."))
                files = [f for f in all_files if not _is_ignored(f)]
                filtered_out = len(all_files) - len(files)
                print(f"Ingesting directory: {p} ({len(files)} files, {filtered_out} skipped by ignore-filter)")
                for f in files:
                    print(f"  Processing: {f.relative_to(p)}")
                    try:
                        count = ingest_file(client, config, collection, f)
                        total += count
                        if count > 0:
                            print(f"    Added {count} chunk(s)")
                        elif count == 0:
                            skipped += 1
                    except Exception as e:
                        print(f"    ERROR: {e}")
                        errors += 1
            elif p.is_file():
                print(f"Ingesting: {p.name}")
                if _is_ignored(p) and not args_force:
                    print(f"  SKIP (ignored): {p}")
                    skipped += 1
                    continue
                try:
                    count = ingest_file(client, config, collection, p)
                    total += count
                    if count > 0:
                        print(f"  Added {count} chunk(s)")
                except Exception as e:
                    print(f"  ERROR: {e}")
                    errors += 1
            else:
                print(f"NOT FOUND: {p}")
    finally:
        _tracker.persist()

    print(f"\nDone! Ingested {total} new chunk(s) into '{collection_name}'")
    if skipped:
        print(f"  Skipped: {skipped} (already existed or empty)")
    if errors:
        print(f"  Errors: {errors}")
    print(_tracker.summary_line())


def cmd_reconcile(args):
    if not args.dry_run and not (args.yes or args.json):
        print("ERROR: Pass --dry-run, --yes, or --json to run reconcile.")
        return

    if args.purge_ignored:
        if args.collection:
            collections = [get_chroma_collection(args.collection)]
        else:
            collections = []
            for item in get_chroma_client().list_collections():
                name = item.name if hasattr(item, "name") else item
                if name.startswith("shared-") or name.startswith("agent-"):
                    collections.append(get_chroma_collection(name))

        report = {
            "dry_run": args.dry_run,
            "purged_files": 0,
            "purged_chunks": 0,
            "kept_files": 0,
            "kept_chunks": 0,
            "reasons": {
                "ignored": {"files": 0, "chunks": 0},
                "missing_from_disk": {"files": 0, "chunks": 0},
            },
            "collections": [],
        }
        for collection in collections:
            collection_report = _purge_collection(collection, dry_run=args.dry_run)
            report["collections"].append(collection_report)
            report["purged_files"] += collection_report["purged_files"]
            report["purged_chunks"] += collection_report["purged_chunks"]
            report["kept_files"] += collection_report["kept_files"]
            report["kept_chunks"] += collection_report["kept_chunks"]
            for reason in report["reasons"]:
                report["reasons"][reason]["files"] += collection_report["reasons"][reason]["files"]
                report["reasons"][reason]["chunks"] += collection_report["reasons"][reason]["chunks"]

        _emit_report(report, json_out=args.json)
        return

    config = load_config()
    client = get_genai_client(get_api_key(config))
    collection_name = args.collection or DEFAULT_RECONCILE_COLLECTION
    collection = get_chroma_collection(collection_name)
    roots = _parse_reconcile_roots(args.roots)
    report = _reconcile_collection(client, config, collection, roots, dry_run=args.dry_run)
    report["orphan_reap"] = _reap_orphan_collections(dry_run=args.dry_run)
    report["reindex"] = _reindex_indexes(DEFAULT_WIKI_ROOT, dry_run=args.dry_run)
    _emit_report(report, json_out=args.json)


def cmd_reindex_indexes(args):
    report = _reindex_indexes(args.root or DEFAULT_WIKI_ROOT, dry_run=args.dry_run)
    if args.json:
        print(json.dumps(report, indent=2))
        return

    print("Reindex summary")
    print("=" * 40)
    print(f"Root: {report['root']}")
    print(f"Dry run: {'yes' if report['dry_run'] else 'no'}")
    print(f"Indexes updated: {report['indexes_updated']}")
    print(f"Indexes created: {report['indexes_created']}")
    print(f"Files linked: {report['files_linked']}")
    if report["dry_run"] and report["diffs"]:
        for diff_text in report["diffs"]:
            print(diff_text, end="" if diff_text.endswith("\n") else "\n")


def cmd_deliver(args):
    config = load_config()
    client = get_genai_client(get_api_key(config))
    collection_name = args.collection or config.get("default_collection", DEFAULT_RECONCILE_COLLECTION)
    collection = get_chroma_collection(collection_name)

    try:
        source_path, resolution = _resolve_delivery_source(
            client,
            config,
            collection,
            args.target,
            top_docs=args.top_docs or DEFAULT_TOP_DOCS,
        )
    except (FileNotFoundError, RuntimeError, ValueError) as exc:
        print(f"ERROR: {exc}")
        return

    if _is_ignored(source_path):
        print(f"ERROR: Refusing to deliver ignored path: {source_path}")
        return
    if not source_path.exists():
        print(f"ERROR: Source file does not exist: {source_path}")
        return

    warning = resolution.get("warning", "")
    if warning and not args.yes:
        print("ERROR: Delivery target is low-confidence or ambiguous. Re-run with --yes to force delivery.")
        return

    try:
        if args.to == "dashboard":
            link = _deliver_to_dashboard(source_path)
        else:
            link = _deliver_to_drive(source_path)
    except RuntimeError as exc:
        print(f"ERROR: {exc}")
        return

    payload = {
        "delivered_path": str(source_path.resolve()),
        "destination": args.to,
        "link": link,
    }
    if args.json:
        print(json.dumps(payload, indent=2))
        return

    print(f"Delivered: {payload['delivered_path']}")
    print(f"Destination: {payload['destination']}")
    print(f"Link: {payload['link']}")


def deduplicate_results(results, similarity_ratio=0.85):
    """Remove near-duplicate results based on content overlap."""
    if len(results) <= 1:
        return results

    deduped = [results[0]]
    for r in results[1:]:
        is_dup = False
        r_content = r["content"][:500]  # compare first 500 chars
        for existing in deduped:
            e_content = existing["content"][:500]
            # Quick overlap check: count shared words
            r_words = set(r_content.lower().split())
            e_words = set(e_content.lower().split())
            if not r_words or not e_words:
                continue
            overlap = len(r_words & e_words) / max(len(r_words), len(e_words))
            if overlap > similarity_ratio:
                is_dup = True
                break
        if not is_dup:
            deduped.append(r)
    return deduped


def _search_query_results(client, config, collection, question, *, top_k=5, threshold=None, type_filter=None, no_recency=False):
    """Run the existing chunk query pipeline and return reranked chunk candidates."""
    fetch_k = (top_k or 5) * 3
    threshold = threshold if threshold is not None else config.get("similarity_threshold", DEFAULT_SIMILARITY_THRESHOLD)

    query_embedding = embed_query(client, config, question)

    where_filter = None
    if type_filter:
        type_map = {
            "image": {"type": "image"},
            "video": {"type": "video_chunk"},
            "text": {"type": "text"},
            "pdf": {"type": "pdf_page"},
            "audio": {"$or": [{"type": "audio"}, {"type": "audio_chunk"}]},
        }
        where_filter = type_map.get(type_filter, {"type": type_filter})

    query_kwargs = {
        "query_embeddings": [query_embedding],
        "n_results": min(fetch_k, collection.count()),
        "include": ["documents", "metadatas", "distances"],
    }
    if where_filter:
        query_kwargs["where"] = where_filter

    try:
        results = collection.query(**query_kwargs)
    except Exception:
        if not where_filter:
            raise
        query_kwargs.pop("where", None)
        results = collection.query(**query_kwargs)

    filtered = []
    if results["ids"] and results["ids"][0]:
        for idx, doc_id in enumerate(results["ids"][0]):
            distance = results["distances"][0][idx] if results["distances"] else 0
            similarity = 1 - distance
            if similarity < threshold:
                continue
            filtered.append({
                "id": doc_id,
                "content": results["documents"][0][idx] if results["documents"] else "",
                "similarity": similarity,
                "metadata": results["metadatas"][0][idx] if results["metadatas"] else {},
            })

    filtered = deduplicate_results(filtered)
    filtered = _apply_recency_rerank(
        filtered,
        enabled=not no_recency,
        weight=config.get("recency_weight", DEFAULT_RECENCY_WEIGHT),
    )
    return filtered, threshold


def _apply_token_budget(results, max_tokens):
    """Trim chunk results to the configured token budget without mutating the originals."""
    if max_tokens <= 0:
        return list(results)

    budgeted = []
    token_count = 0
    for result in results:
        chunk_tokens = len(result["content"]) // 4
        if token_count + chunk_tokens > max_tokens:
            remaining = max_tokens - token_count
            if remaining > 50:
                budgeted.append({
                    **result,
                    "content": result["content"][:remaining * 4] + "... [truncated]",
                })
            break
        budgeted.append(result)
        token_count += chunk_tokens
    return budgeted


def _result_score(result):
    return result.get("final_score", result["similarity"])


def _group_parent_documents(results, *, top_docs=None):
    """Group chunk hits by source document and score each parent doc."""
    documents = {}
    for result in results:
        source = _metadata_source(result.get("metadata"))
        if not source:
            continue
        entry = documents.setdefault(source, {
            "source_file": source,
            "abs_path": str(Path(source).resolve()),
            "matched_chunks": 0,
            "doc_score": float("-inf"),
            "best_similarity": 0.0,
            "title": _humanize_filename(Path(source).stem),
            "highlight": "",
        })
        entry["matched_chunks"] += 1
        entry["best_similarity"] = max(entry["best_similarity"], result["similarity"])
        score = _result_score(result)
        if score > entry["doc_score"]:
            title, _ = _document_title_and_summary(Path(source))
            entry["title"] = title
            entry["doc_score"] = score
            entry["highlight"] = result["content"][:DEFAULT_PREVIEW_CHARS]

    grouped = sorted(
        (
            {
                **entry,
                "doc_score": round(entry["doc_score"], 4),
                "best_similarity": round(entry["best_similarity"], 4),
            }
            for entry in documents.values()
        ),
        key=lambda entry: (-entry["doc_score"], -entry["matched_chunks"], entry["source_file"]),
    )
    if top_docs is None:
        return grouped
    return grouped[:top_docs]


def _read_full_document(file_path: Path, *, max_bytes=MAX_FULL_DOC_BYTES):
    """Read a source document from disk for --docs --full delivery."""
    file_path = Path(file_path).expanduser().resolve()
    if not file_path.exists():
        raise FileNotFoundError(str(file_path))
    if _is_ignored(file_path):
        raise ValueError(f"Ignored path cannot be returned as a full document: {file_path}")
    content, truncated = _read_text_file(file_path, max_bytes=max_bytes)
    result = {
        "source_file": str(file_path),
        "abs_path": str(file_path),
        "content": content,
        "truncated": truncated,
    }
    if truncated:
        result["note"] = f"Document exceeded {max_bytes} bytes; returning the first {max_bytes} bytes."
    return result


def _build_chunk_output_entry(result, *, show_full=False):
    """Stable JSON chunk entry used by both legacy and doc-grain output."""
    metadata = result["metadata"]
    source = _metadata_source(metadata)
    entry = {
        "content": result["content"] if show_full else result["content"][:DEFAULT_PREVIEW_CHARS],
        "content_full_length": len(result["content"]),
        "similarity": round(result["similarity"], 4),
        "source": source,
        "source_file": source,
        "type": metadata.get("type", ""),
        "filename": metadata.get("filename", ""),
    }
    if metadata.get("chunk_index") is not None:
        entry["chunk_index"] = metadata["chunk_index"]
        entry["total_chunks"] = metadata.get("total_chunks", 0)
    if metadata.get("chunk_start_seconds") is not None:
        entry["time_start"] = metadata["chunk_start_seconds"]
        entry["time_end"] = metadata.get("chunk_end_seconds", 0)
    if metadata.get("chunk_path"):
        entry["chunk_path"] = metadata["chunk_path"]
    if metadata.get("page_number"):
        entry["page_number"] = metadata["page_number"]
    return entry


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[2]


def _fleet_agents_dir() -> Path:
    org_name = os.environ.get("CTX_ORG", "clearworksai")
    return _repo_root() / "orgs" / org_name / "agents"


def _fleet_agent_names():
    agents_dir = _fleet_agents_dir()
    if not agents_dir.exists():
        return set()
    return {
        path.name for path in agents_dir.iterdir()
        if path.is_dir()
    }


def _collection_owner_agent(collection_name: str, agent_names) -> str:
    if not collection_name.startswith("agent-"):
        return ""
    suffix = collection_name[len("agent-"):]
    if suffix in agent_names:
        return suffix
    for agent_name in sorted(agent_names, key=len, reverse=True):
        if suffix.startswith(f"{agent_name}-") or suffix.endswith(f"-{agent_name}"):
            return agent_name
    return ""


def _reap_orphan_collections(chroma=None, *, dry_run=False, agent_names=None):
    """Delete empty agent-scoped collections and report anything suspicious we kept."""
    chroma = chroma or get_chroma_client()
    agent_names = set(agent_names or _fleet_agent_names())
    report = {
        "dry_run": dry_run,
        "collections_scanned": 0,
        "orphans_found": 0,
        "orphans_reaped": 0,
        "kept": 0,
        "orphan_names": [],
        "kept_names": [],
        "suspicious_nonempty": [],
    }

    for item in chroma.list_collections():
        name = item.name if hasattr(item, "name") else item
        report["collections_scanned"] += 1
        if name.startswith("shared-") or not name.startswith("agent-"):
            report["kept"] += 1
            report["kept_names"].append(name)
            continue

        collection = chroma.get_collection(name)
        count = collection.count()
        owner = _collection_owner_agent(name, agent_names)
        if count == 0:
            report["orphans_found"] += 1
            report["orphans_reaped"] += 1
            report["orphan_names"].append(name)
            if not dry_run:
                chroma.delete_collection(name)
            continue

        if not owner:
            report["suspicious_nonempty"].append(name)
        report["kept"] += 1
        report["kept_names"].append(name)

    return report


def _run_command(command, *, cwd=None, env=None, input_text=None):
    """Run a subprocess and return stdout, raising a readable error on failure."""
    try:
        result = subprocess.run(
            command,
            cwd=cwd,
            env=env,
            input=input_text,
            capture_output=True,
            text=True,
            check=False,
        )
    except OSError as exc:
        raise RuntimeError(str(exc)) from exc
    if result.returncode != 0:
        stderr = result.stderr.strip() or result.stdout.strip() or f"Command failed: {' '.join(command)}"
        raise RuntimeError(stderr)
    return result.stdout.strip()


def _extract_link_from_payload(payload):
    """Find a URL-like field in a nested CLI response payload."""
    if isinstance(payload, dict):
        for key in ("link", "url", "webViewLink", "web_url", "alternateLink"):
            value = payload.get(key)
            if isinstance(value, str) and value.startswith("http"):
                return value
        for value in payload.values():
            link = _extract_link_from_payload(value)
            if link:
                return link
        return ""
    if isinstance(payload, list):
        for value in payload:
            link = _extract_link_from_payload(value)
            if link:
                return link
    return ""


def _build_dashboard_delivery_content(file_path: Path) -> str:
    """Render a source file into markdown for the briefs delivery surface."""
    document = _read_full_document(file_path)
    lines = [
        f"## {file_path.name}",
        f"Source path: `{file_path}`",
        "",
        "```text",
        document["content"],
        "```",
    ]
    if document.get("note"):
        lines.extend(["", document["note"]])
    return "\n".join(lines)


def _deliver_to_dashboard(file_path: Path) -> str:
    """Publish a source file into the briefs/dashboard surface and return its URL."""
    briefs_repo = _repo_root().parent / "briefs"
    publisher = briefs_repo / "publisher" / "publish_brief.py"
    title = f"Delivered File - {file_path.name}"
    return _run_command(
        [sys.executable, str(publisher), "--title", title],
        cwd=str(briefs_repo),
        input_text=_build_dashboard_delivery_content(file_path),
    )


def _deliver_to_drive(file_path: Path) -> str:
    """Upload a file to Josh's Drive surface through the existing CLI path."""
    candidates = []
    if shutil.which("gog"):
        candidates.extend([
            ["gog", "drive", "upload", "-a", "josh@clearworks.ai", "--json", str(file_path)],
            ["gog", "drive", "put", "-a", "josh@clearworks.ai", "--json", str(file_path)],
        ])
    if shutil.which("officecli"):
        candidates.append(["officecli", "drive", "upload", str(file_path), "--json"])
    if not candidates:
        raise RuntimeError("No Drive delivery CLI is available in PATH")

    last_error = None
    for command in candidates:
        try:
            output = _run_command(command)
        except RuntimeError as exc:
            last_error = exc
            continue
        if output.startswith("http"):
            return output
        try:
            payload = json.loads(output)
        except json.JSONDecodeError:
            last_error = RuntimeError(f"Drive delivery returned non-JSON output: {output}")
            continue
        link = _extract_link_from_payload(payload)
        if link:
            return link
        last_error = RuntimeError("Drive delivery completed without a shareable link")

    raise last_error or RuntimeError("Drive delivery failed")


def _resolve_delivery_source(client, config, collection, target: str, *, top_docs=DEFAULT_TOP_DOCS, no_recency=False):
    """Resolve a path or free-text query into a concrete source file."""
    candidate = Path(target).expanduser()
    if candidate.exists():
        resolved = candidate.resolve()
        return resolved, {"mode": "path", "documents": []}

    chunk_results, _ = _search_query_results(
        client,
        config,
        collection,
        target,
        top_k=max(top_docs, 5),
        threshold=DEFAULT_SIMILARITY_THRESHOLD,
        type_filter=None,
        no_recency=no_recency,
    )
    documents = _group_parent_documents(chunk_results, top_docs=max(top_docs, 3))
    if not documents:
        raise FileNotFoundError(f"No documents matched query: {target}")

    warning = ""
    if documents[0]["doc_score"] < DELIVER_LOW_SCORE:
        warning = "low_score"
    elif len(documents) > 1 and documents[0]["doc_score"] - documents[1]["doc_score"] < DELIVER_AMBIGUITY_GAP:
        warning = "ambiguous"

    return Path(documents[0]["abs_path"]), {
        "mode": "query",
        "documents": documents,
        "warning": warning,
    }


def cmd_query(args):
    global _tracker
    _tracker = UsageTracker("query")

    config = load_config()
    client = get_genai_client(get_api_key(config))
    collection_name = args.collection or config.get("default_collection", "default")
    collection = get_chroma_collection(collection_name)

    if collection.count() == 0:
        print("Knowledge base is empty. Ingest some files first.")
        return

    max_tokens = args.max_tokens or config.get("max_tokens", DEFAULT_MAX_TOKENS)
    docs_mode = getattr(args, "docs", False) or getattr(args, "parent", False)

    chunk_candidates, threshold = _search_query_results(
        client,
        config,
        collection,
        args.question,
        top_k=args.top_k or 5,
        threshold=args.threshold,
        type_filter=args.type,
        no_recency=getattr(args, "no_recency", False),
    )
    document_limit = getattr(args, "top_docs", None) or (args.top_k or DEFAULT_TOP_DOCS)
    document_results = _group_parent_documents(chunk_candidates, top_docs=document_limit)

    final_k = args.top_k or 5
    filtered = _apply_token_budget(chunk_candidates[:final_k], max_tokens)

    # Collect unique source files for the agent
    source_files = list(dict.fromkeys(
        _metadata_source(r["metadata"]) for r in filtered if _metadata_source(r["metadata"])
    ))

    full_document = None
    if docs_mode and args.full and document_results:
        try:
            full_document = _read_full_document(Path(document_results[0]["abs_path"]))
        except (OSError, ValueError) as exc:
            full_document = {
                "source_file": document_results[0]["source_file"],
                "abs_path": document_results[0]["abs_path"],
                "error": str(exc),
            }

    if args.json:
        chunk_entries = [_build_chunk_output_entry(result, show_full=args.full) for result in filtered]
        output = {
            "query": args.question,
            "collection": collection_name,
            "result_count": len(filtered),
            "source_files": source_files,
            "results": chunk_entries,
            "chunks": chunk_entries,
            "documents": document_results,
        }
        if full_document is not None:
            output["full_document"] = full_document

        print(json.dumps(output, indent=2))
    else:
        if docs_mode:
            print(f"Query: {args.question}")
            print(f"Collection: {collection_name}")
            print(f"Documents: {len(document_results)} (threshold: {threshold})")
            print("-" * 60)
            if document_results:
                for idx, document in enumerate(document_results, start=1):
                    print(f"\n[{idx}] Score: {document['doc_score']:.3f}")
                    print(f"    Title: {document['title']}")
                    print(f"    Source: {document['source_file']}")
                    print(f"    Matched chunks: {document['matched_chunks']}")
                    print(f"    Path: {document['abs_path']}")
                    if document["highlight"]:
                        print(f"    Highlight: {document['highlight']}")
                if full_document is not None:
                    print("\nFull document")
                    print("-" * 60)
                    if full_document.get("error"):
                        print(full_document["error"])
                    else:
                        print(full_document["content"])
                        if full_document.get("note"):
                            print()
                            print(full_document["note"])
            else:
                print("No documents above similarity threshold.")
        else:
            print(f"Query: {args.question}")
            print(f"Collection: {collection_name}")
            print(f"Results: {len(filtered)} (threshold: {threshold})")
            if source_files:
                print(f"Source files ({len(source_files)}):")
                for sf in source_files:
                    print(f"  - {sf}")
            print("-" * 60)

            if filtered:
                for i, result in enumerate(filtered):
                    metadata = result["metadata"]
                    print(f"\n[{i+1}] Similarity: {result['similarity']:.3f}")
                    print(f"    Source: {_metadata_source(metadata) or 'unknown'}")
                    print(f"    Type: {metadata.get('type', 'unknown')}")
                    if metadata.get("chunk_index") is not None:
                        print(f"    Chunk: {metadata['chunk_index'] + 1}/{metadata.get('total_chunks', '?')}")
                    if metadata.get("chunk_start_seconds") is not None:
                        print(f"    Time: {metadata['chunk_start_seconds']:.0f}s - {metadata.get('chunk_end_seconds', 0):.0f}s")
                    if metadata.get("chunk_path"):
                        print(f"    Chunk file: {metadata['chunk_path']}")
                    if metadata.get("page_number"):
                        print(f"    Page: {metadata['page_number']}")

                    content = result["content"]
                    if not args.full and len(content) > DEFAULT_PREVIEW_CHARS:
                        content = content[:DEFAULT_PREVIEW_CHARS] + f"... [{len(result['content'])} chars total, use --full to see all]"
                    print(f"    Content: {content}")
            else:
                print("No results above similarity threshold.")

    _tracker.persist()


def cmd_usage(args):
    """Show token usage and cost summary."""
    if args.reset:
        if USAGE_FILE.exists():
            USAGE_FILE.unlink()
        print("Usage data reset.")
        return

    if not USAGE_FILE.exists():
        print("No usage data yet. Run an ingest or query first.")
        return

    with open(USAGE_FILE) as f:
        data = json.load(f)

    if args.json:
        print(json.dumps(data, indent=2))
        return

    c = data.get("cumulative", {})
    sessions = data.get("sessions", [])

    emb_cost = (c.get("embedding_tokens", 0) / 1_000_000) * EMBEDDING_PRICE_PER_M
    gen_in_cost = (c.get("generation_input_tokens", 0) / 1_000_000) * FLASH_INPUT_PRICE_PER_M
    gen_out_cost = (c.get("generation_output_tokens", 0) / 1_000_000) * FLASH_OUTPUT_PRICE_PER_M
    total = emb_cost + gen_in_cost + gen_out_cost

    print("mmrag Usage Summary")
    print("=" * 40)
    print(f"Total cost:    ${total:.4f}")
    print(f"Total calls:   {c.get('embedding_calls', 0) + c.get('generation_calls', 0)} "
          f"({c.get('embedding_calls', 0)} embedding, {c.get('generation_calls', 0)} generation)")
    print()
    print("Token breakdown:")
    print(f"  Embedding:         {c.get('embedding_tokens', 0):>10,} tokens (est)  ${emb_cost:.4f}")
    print(f"  Generation input:  {c.get('generation_input_tokens', 0):>10,} tokens        ${gen_in_cost:.4f}")
    print(f"  Generation output: {c.get('generation_output_tokens', 0):>10,} tokens        ${gen_out_cost:.4f}")
    print()
    print(f"Sessions: {len(sessions)}")
    if sessions:
        last = sessions[-1]
        print(f"  Last: {last.get('operation', '?')} @ {last.get('finished_at', '?')} "
              f"(${last.get('cost', {}).get('total', 0):.4f})")

    # Per-day breakdown from sessions
    by_day = {}
    for s in sessions:
        day = s.get("started_at", "")[:10]
        if day:
            by_day.setdefault(day, {"cost": 0, "count": 0})
            by_day[day]["cost"] += s.get("cost", {}).get("total", 0)
            by_day[day]["count"] += 1

    if by_day:
        print()
        print("Daily breakdown:")
        for day in sorted(by_day.keys(), reverse=True)[:7]:
            d = by_day[day]
            print(f"  {day}:  ${d['cost']:.4f} ({d['count']} sessions)")


def cmd_status(args):
    config = load_config()
    collection_name = args.collection or config.get("default_collection", "default")

    try:
        collection = get_chroma_collection(collection_name)
        count = collection.count()
    except Exception:
        count = 0

    print(f"Collection: {collection_name}")
    print(f"Total chunks: {count}")
    print(f"Data dir: {MMRAG_DIR}")
    print(f"ChromaDB: {CHROMADB_DIR}")
    print(f"Config: {CONFIG_FILE}")

    # Show config values
    print(f"\nChunk settings:")
    print(f"  Text: {config.get('text_chunk_size', DEFAULT_TEXT_CHUNK_SIZE)} chars, {config.get('text_chunk_overlap', DEFAULT_TEXT_CHUNK_OVERLAP)} overlap")
    print(f"  Video: {config.get('video_chunk_seconds', DEFAULT_VIDEO_CHUNK_SECONDS)}s, {config.get('video_overlap_seconds', DEFAULT_VIDEO_OVERLAP_SECONDS)}s overlap")
    print(f"  Audio: {config.get('audio_chunk_seconds', DEFAULT_AUDIO_CHUNK_SECONDS)}s, {config.get('audio_overlap_seconds', DEFAULT_AUDIO_OVERLAP_SECONDS)}s overlap")
    print(f"  Embedding dims: {config.get('embedding_dimensions', DEFAULT_EMBEDDING_DIMENSIONS)}")

    if count > 0:
        all_data = _get_all_metadatas(collection)
        types_map = {}
        sources = set()
        for meta in all_data["metadatas"]:
            t = meta.get("type", "unknown")
            types_map[t] = types_map.get(t, 0) + 1
            sources.add(_metadata_source(meta) or "unknown")

        print(f"\nUnique sources: {len(sources)}")
        print("Type breakdown:")
        for t, c in sorted(types_map.items()):
            print(f"  {t}: {c} chunks")


def cmd_list(args):
    config = load_config()
    collection_name = args.collection or config.get("default_collection", "default")

    try:
        collection = get_chroma_collection(collection_name)
    except Exception:
        print("No data found.")
        return

    all_data = _get_all_metadatas(collection)
    if not all_data["ids"]:
        print("No documents in collection.")
        return

    # Group by source
    by_source = {}
    for meta in all_data["metadatas"]:
        src = _metadata_source(meta) or "unknown"
        if src not in by_source:
            by_source[src] = {"type": meta.get("type", "unknown"), "chunks": 0,
                              "filename": meta.get("filename", "")}
        by_source[src]["chunks"] += 1

    print(f"Collection: {collection_name} ({len(by_source)} files, {len(all_data['ids'])} chunks)")
    print(f"{'Source':<60} {'Type':<15} {'Chunks':<8}")
    print("-" * 85)
    for src, info in sorted(by_source.items()):
        display = src if len(src) <= 58 else "..." + src[-55:]
        print(f"{display:<60} {info['type']:<15} {info['chunks']:<8}")


def cmd_collections(args):
    chroma = get_chroma_client()
    collections = chroma.list_collections()
    if not collections:
        print("No collections found.")
        return
    print(f"{'Collection':<30} {'Documents':<12}")
    print("-" * 44)
    for c in collections:
        col = chroma.get_collection(c.name if hasattr(c, 'name') else c)
        name = c.name if hasattr(c, 'name') else c
        print(f"{name:<30} {col.count():<12}")


def cmd_delete(args):
    config = load_config()
    collection_name = args.collection or config.get("default_collection", "default")
    collection = get_chroma_collection(collection_name)

    source_path = str(Path(args.path).resolve())
    all_data = _get_all_metadatas(collection)

    ids_to_delete = []
    for i, meta in enumerate(all_data["metadatas"]):
        if _metadata_source(meta) == source_path:
            ids_to_delete.append(all_data["ids"][i])

    if not ids_to_delete:
        print(f"No documents found for: {source_path}")
        return

    collection.delete(ids=ids_to_delete)
    print(f"Deleted {len(ids_to_delete)} chunk(s) from '{collection_name}' for: {source_path}")


def cmd_reset(args):
    if not args.confirm:
        print("ERROR: Pass --confirm to reset the knowledge base.")
        return

    chroma = get_chroma_client()
    collections = chroma.list_collections()
    count = 0
    for c in collections:
        name = c.name if hasattr(c, 'name') else c
        chroma.delete_collection(name)
        count += 1

    import shutil
    if MEDIA_DIR.exists():
        shutil.rmtree(MEDIA_DIR)
        MEDIA_DIR.mkdir(parents=True)

    print(f"Reset complete. Deleted {count} collection(s) and cleared media cache.")

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    parser = argparse.ArgumentParser(
        description="Multimodal RAG Knowledge Base CLI",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    sub = parser.add_subparsers(dest="command", help="Command to run")

    # ingest
    p_ingest = sub.add_parser("ingest", help="Ingest files into the knowledge base")
    p_ingest.add_argument("paths", nargs="+", help="File or directory paths to ingest")
    p_ingest.add_argument("--collection", "-c", help="Collection name (default: 'default')")
    p_ingest.add_argument("--force", action="store_true", help="Re-ingest files even if already in the KB")

    # query
    p_query = sub.add_parser("query", help="Query the knowledge base")
    p_query.add_argument("question", help="Question to ask")
    p_query.add_argument("--top-k", "-k", type=int, default=5, help="Max number of results (default: 5)")
    p_query.add_argument("--threshold", "-t", type=float, default=None,
                         help="Min similarity threshold 0.0-1.0 (default: 0.0, return all)")
    p_query.add_argument("--max-tokens", "-m", type=int, default=0,
                         help="Max total tokens in results (0=unlimited)")
    p_query.add_argument("--collection", "-c", help="Collection name")
    p_query.add_argument("--type", help="Filter by content type: image, video, text, pdf, audio")
    p_query.add_argument("--json", "-j", action="store_true", help="Output as JSON (for agent consumption)")
    p_query.add_argument("--full", "-f", action="store_true", help="Show full content (not truncated)")
    p_query.add_argument("--docs", action="store_true", help="Return document-grain results grouped by source_file")
    p_query.add_argument("--parent", action="store_true", help="Alias for --docs")
    p_query.add_argument("--top-docs", type=int, default=DEFAULT_TOP_DOCS,
                         help=f"Max number of parent documents to return in --docs mode (default: {DEFAULT_TOP_DOCS})")
    p_query.add_argument("--no-recency", action="store_true",
                         help="Disable recency reranking and preserve pure similarity ordering")

    # status
    p_status = sub.add_parser("status", help="Show knowledge base status")
    p_status.add_argument("--collection", "-c", help="Collection name")

    # list
    p_list = sub.add_parser("list", help="List ingested documents")
    p_list.add_argument("--collection", "-c", help="Collection name")

    # collections
    sub.add_parser("collections", help="List all collections")

    # delete
    p_delete = sub.add_parser("delete", help="Delete a document by source path")
    p_delete.add_argument("path", help="Source file path to delete")
    p_delete.add_argument("--collection", "-c", help="Collection name")

    # reset
    p_reset = sub.add_parser("reset", help="Reset the entire knowledge base")
    p_reset.add_argument("--confirm", action="store_true", help="Confirm reset")

    # reconcile
    p_reconcile = sub.add_parser("reconcile", help="Reconcile the indexed corpus against disk")
    p_reconcile.add_argument("--collection", "-c",
                             help="Collection name (default: shared-clearworksai, or all shared/agent collections with --purge-ignored)")
    p_reconcile.add_argument("--roots",
                             help="Comma-separated roots to walk (default: ~/code/knowledge-sync/wiki,~/code/knowledge-sync/raw)")
    p_reconcile.add_argument("--dry-run", action="store_true", help="Show what would change without mutating the collection")
    p_reconcile.add_argument("--json", "-j", action="store_true", help="Output reconcile results as JSON")
    p_reconcile.add_argument("--yes", action="store_true", help="Confirm destructive reconcile changes")
    p_reconcile.add_argument("--purge-ignored", action="store_true",
                             help="Only purge indexed chunks whose sources are ignored or missing from disk")
    p_reconcile.add_argument("--reap-orphans", action="store_true",
                             help="Also delete empty orphaned agent-* collections during reconcile")

    # reindex-indexes
    p_reindex = sub.add_parser("reindex-indexes", help="Refresh managed _index.md sections under the wiki root")
    p_reindex.add_argument("--root", help="Wiki root to scan (default: ~/code/knowledge-sync/wiki)")
    p_reindex.add_argument("--dry-run", action="store_true", help="Show diffs without writing any _index.md files")
    p_reindex.add_argument("--json", "-j", action="store_true", help="Output reindex results as JSON")

    # deliver
    p_deliver = sub.add_parser("deliver", help="Resolve and deliver a source file to Josh's Drive or dashboard")
    p_deliver.add_argument("target", help="Exact path or free-text query to resolve to a source file")
    p_deliver.add_argument("--to", choices=("drive", "dashboard"), default="drive",
                           help="Delivery surface (default: drive)")
    p_deliver.add_argument("--collection", "-c", help="Collection name (default: shared-clearworksai)")
    p_deliver.add_argument("--top-docs", type=int, default=DEFAULT_TOP_DOCS,
                           help=f"How many documents to inspect when resolving a query (default: {DEFAULT_TOP_DOCS})")
    p_deliver.add_argument("--yes", action="store_true",
                           help="Force delivery even when query resolution is low-confidence or ambiguous")
    p_deliver.add_argument("--json", "-j", action="store_true", help="Output delivery results as JSON")

    # usage
    p_usage = sub.add_parser("usage", help="Show token usage and cost summary")
    p_usage.add_argument("--json", "-j", action="store_true", help="Output as JSON")
    p_usage.add_argument("--reset", action="store_true", help="Reset usage data")

    args = parser.parse_args()

    if not args.command:
        parser.print_help()
        sys.exit(1)

    commands = {
        "ingest": cmd_ingest,
        "query": cmd_query,
        "status": cmd_status,
        "list": cmd_list,
        "collections": cmd_collections,
        "delete": cmd_delete,
        "reset": cmd_reset,
        "reconcile": cmd_reconcile,
        "reindex-indexes": cmd_reindex_indexes,
        "deliver": cmd_deliver,
        "usage": cmd_usage,
    }

    commands[args.command](args)


if __name__ == "__main__":
    main()
